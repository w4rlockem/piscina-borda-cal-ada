"""Monta o PPTX de 18 slides a partir das imagens em deck/render/."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import montagem as M
from gerar import SAIDA, COMBOS, BORDA_ALTERNATIVA

DESTINO = M.caminho("Piscina menor - materiais de borda e calcada.pptx")

L, A = 13.333, 7.5                       # slide 16:9, em polegadas

TINTA = RGBColor(0x1A, 0x1A, 0x1C)
SUAVE = RGBColor(0x70, 0x70, 0x78)
ACENTO = RGBColor(0x1E, 0x6F, 0x84)
FUNDO = RGBColor(0xFF, 0xFF, 0xFF)
CREME = RGBColor(0xF4, 0xF2, 0xEE)
BORDA_C = RGBColor(0xDD, 0xDA, 0xD4)
ALERTA = RGBColor(0xA8, 0x38, 0x2A)
VERDE = RGBColor(0x3F, 0x6B, 0x46)

TIT = "Georgia"
TXT = "Calibri"

CRITERIOS = [
    ("aderencia", "Antiderrapante"),
    ("termico", "Não esquenta"),
    ("custo", "Custo baixo"),
    ("manutencao", "Baixa manutenção"),
    ("maresia", "Resiste à maresia"),
]

SELO = {"sim": ("Serve", VERDE), "limitado": ("Limitado", RGBColor(0xB5, 0x7A, 0x1F)),
        "nao": ("Não serve", ALERTA)}


# ---------------------------------------------------------------- utilidades


def novo(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def caixa(slide, x, y, w, h, txt, tam=12, cor=TINTA, neg=False, fam=TXT,
          align=PP_ALIGN.LEFT, entre=1.0, italico=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, linha in enumerate(str(txt).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = entre
        r = p.add_run()
        r.text = linha
        r.font.size = Pt(tam)
        r.font.color.rgb = cor
        r.font.bold = neg
        r.font.italic = italico
        r.font.name = fam
    return tb


def bloco(slide, x, y, w, h, cor, linha=None, esp=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = cor
    if linha is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = linha
        s.line.width = Pt(esp)
    s.shadow.inherit = False
    return s


def figura(slide, arquivo, x, y, w):
    return slide.shapes.add_picture(os.path.join(SAIDA, arquivo),
                                    Inches(x), Inches(y), width=Inches(w))


def pips(slide, x, y, nota, largura=0.17, alt=0.115, gap=0.035):
    """Cinco quadradinhos preenchidos conforme a nota de 1 a 5."""
    for i in range(5):
        cor = ACENTO if i < nota else RGBColor(0xE2, 0xE0, 0xDB)
        bloco(slide, x + i * (largura + gap), y, largura, alt, cor)


def faixa_preco(mat):
    meio = sum(mat["preco_m2"]) / 2
    return "$" * (1 if meio < 200 else 2 if meio < 280 else 3 if meio < 380 else 4)


def rodape(slide, texto_):
    caixa(slide, 0.55, A - 0.42, 12.2, 0.3, texto_, 8.5, SUAVE)


def salva(prs, destino):
    """Salva o PPTX. Se o arquivo estiver aberto no PowerPoint o Windows
    trava a escrita -- nesse caso grava ao lado com sufixo, em vez de
    perder a geracao inteira."""
    try:
        prs.save(destino)
        return destino
    except PermissionError:
        base, ext = os.path.splitext(destino)
        alt = f"{base} (novo){ext}"
        prs.save(alt)
        print("AVISO: o arquivo original esta aberto no PowerPoint.")
        print("       Feche-o e renomeie, ou rode de novo com ele fechado.")
        return alt


# ---------------------------------------------------------------- slides


def s01_capa(prs):
    s = novo(prs)
    figura(s, "capa.jpg", 0, -0.55, L)
    bloco(s, 0, 5.35, L, 2.15, RGBColor(0x12, 0x2A, 0x33))
    caixa(s, 0.85, 5.62, 11, 0.6, "Piscina menor — troca de borda e calçada",
          30, RGBColor(0xFF, 0xFF, 0xFF), True, TIT)
    caixa(s, 0.85, 6.28, 11, 0.4,
          "10 materiais aplicados sobre as fotos do local",
          15, RGBColor(0xC9, 0xDA, 0xE0), False, TXT)
    caixa(s, 0.85, 6.78, 11, 0.35,
          "Ponta da Fruta, Vila Velha — ES  ·  agosto de 2026",
          11, RGBColor(0x91, 0xAA, 0xB4))


def s02_diagnostico(prs):
    s = novo(prs)
    caixa(s, 0.55, 0.42, 8, 0.5, "O que existe hoje", 26, TINTA, True, TIT)
    figura(s, "diagnostico.jpg", 0.55, 1.15, 7.1)
    caixa(s, 0.55, 6.55, 7.1, 0.3,
          "Detalhe do canto: a fibra termina e sobra concreto bruto até a pedra.",
          9, SUAVE, italico=True)

    x, y = 8.15, 1.15
    itens = [
        ("1", "Não existe peça de borda",
         "A casca de fibra termina e sobra uma faixa de concreto bruto até a "
         "pedra. É o defeito principal e responde pela maior parte da sensação "
         "de malacabado."),
        ("2", "Rejunte largo, manchado, com mato",
         "Degrada mais a percepção da área do que o desgaste da pedra em si."),
        ("3", "Duas paginações misturadas",
         "Placa serrada regular na maior parte e caco irregular junto à "
         "piscina, sem transição."),
    ]
    for num, tit, desc in itens:
        bloco(s, x, y + 0.03, 0.30, 0.30, ALERTA)
        caixa(s, x + 0.085, y + 0.06, 0.3, 0.3, num, 12,
              RGBColor(0xFF, 0xFF, 0xFF), True)
        caixa(s, x + 0.45, y + 0.02, 4.3, 0.3, tit, 13, TINTA, True)
        caixa(s, x + 0.45, y + 0.36, 4.3, 1.0, desc, 10.5, SUAVE, entre=1.15)
        y += 1.42

    bloco(s, x, 5.62, 4.65, 1.35, CREME)
    caixa(s, x + 0.25, 5.8, 4.2, 0.25, "LEVANTAMENTO ESTIMADO", 9, ACENTO, True)
    caixa(s, x + 0.25, 6.12, 4.2, 0.8,
          "Piscina ~6,0 × 3,0 m  ·  calçada de 1,5 a 2,5 m no perímetro\n"
          "Piso: 60 a 75 m²  ·  Borda: ~18 m lineares", 11, TINTA, entre=1.3)
    caixa(s, x + 0.25, 6.72, 4.2, 0.25,
          "Estimado por foto — confirmar com trena.", 8.5, SUAVE, italico=True)


def s03_o_que_muda(prs):
    s = novo(prs)
    caixa(s, 0.55, 0.40, 9, 0.45, "O que muda", 26, TINTA, True, TIT)
    caixa(s, 0.55, 0.92, 12.2, 0.3,
          "Sai a São Tomé inteira. Entram borda nova e piso novo, sobre "
          "contrapiso regularizado.", 12, SUAVE)
    figura(s, "perfis.png", 0.95, 1.28, 11.45)

    caixa(s, 0.55, 4.92, 12.2, 0.3,
          "Seis detalhes que separam obra boa de obra que estraga em dois anos",
          14, TINTA, True, TIT)
    detalhes = [
        ("A borda avança 2–3 cm sobre a fibra, com pingadeira",
         "Sem o sulco, a água escorre pela face e mancha."),
        ("Junta flexível na interface, nunca rejunte rígido",
         "A fibra trabalha com temperatura; rejunte duro trinca."),
        ("Caimento de 1 a 1,5% para fora da piscina",
         "Hoje a água suja da calçada volta para dentro."),
        ("Rejunte epóxi, não cimentício",
         "Resiste ao cloro, não mancha e não cria limo."),
        ("Juntas de dilatação a cada ~3 m",
         "Em 60–75 m² de sol pleno, sem elas o piso estufa."),
        ("Ferragem em inox 316 (A4), não A2",
         "Litoral: A2 mancha de ferrugem no primeiro ano."),
    ]
    for i, (tit, desc) in enumerate(detalhes):
        cx = 0.55 + (i % 3) * 4.13
        cy = 5.35 + (i // 3) * 1.06
        bloco(s, cx, cy, 3.88, 0.95, CREME)
        bloco(s, cx, cy, 0.045, 0.95, ACENTO)
        caixa(s, cx + 0.22, cy + 0.12, 3.5, 0.42, tit, 10.5, TINTA, True,
              entre=1.1)
        caixa(s, cx + 0.22, cy + 0.58, 3.5, 0.32, desc, 9, SUAVE, entre=1.1)


CRIT_BORDA = [
    ("aderencia", "Antiderrapante"),
    ("termico", "Não esquenta"),
    ("cloro", "Resiste ao cloro"),
    ("sob_medida", "Peça sob medida"),
    ("custo", "Custo baixo"),
]


def s_borda_intro(prs):
    s = novo(prs)
    caixa(s, 0.55, 0.38, 10, 0.45, "A borda: o que realmente muda", 26, TINTA,
          True, TIT)
    caixa(s, 0.55, 0.93, 12.2, 0.3,
          "A peça avança 2–3 cm sobre a lâmina da fibra e a esconde. É a "
          "diferença entre piscina acabada e piscina montada.", 12, SUAVE)

    for i, (arq, rot) in enumerate([
            ("borda_antes.jpg", "HOJE — a borda é a própria fibra"),
            ("borda_b-granito-branco.jpg", "COM PEÇA DE BORDA em granito")]):
        cx = 0.55 + i * 6.24
        figura(s, arq, cx, 1.42, 5.99)
        cor = ALERTA if i == 0 else VERDE
        bloco(s, cx, 5.05, 5.99, 0.32, cor)
        caixa(s, cx + 0.15, 5.11, 5.7, 0.28, rot, 10,
              RGBColor(0xFF, 0xFF, 0xFF), True)

    itens = [("A peça esconde a lâmina de fibra",
              "Hoje o rebordo azul fica à mostra e sobra concreto bruto "
              "até a pedra."),
             ("A borda é comprada por metro linear",
              "São ~18 m no perímetro desta piscina, e a peça é cortada "
              "sob medida."),
             ("Pode ser de material diferente do piso",
              "É comum e recomendável: a borda pede peça boleada com "
              "pingadeira.")]
    for i, (tit, desc) in enumerate(itens):
        cx = 0.55 + i * 4.13
        bloco(s, cx, 5.62, 3.88, 1.35, CREME)
        bloco(s, cx, 5.62, 0.045, 1.35, ACENTO)
        caixa(s, cx + 0.22, 5.78, 3.5, 0.4, tit, 11, TINTA, True, entre=1.1)
        caixa(s, cx + 0.22, 6.28, 3.5, 0.6, desc, 9.5, SUAVE, entre=1.15)


def s_borda(prs, b):
    s = novo(prs)
    bloco(s, 0, 0, L, 0.92, CREME)
    bloco(s, 0, 0.92, L, 0.02, ACENTO)
    caixa(s, 0.55, 0.16, 0.7, 0.5, str(b["n"]), 26, ACENTO, True, TIT)
    caixa(s, 1.25, 0.13, 7.6, 0.4, b["nome"], 21, TINTA, True, TIT)
    caixa(s, 1.25, 0.575, 7.6, 0.28, b["subtitulo"], 10.5, SUAVE)
    caixa(s, 9.3, 0.17, 3.5, 0.3,
          f"R$ {b['preco_ml'][0]}–{b['preco_ml'][1]} / m linear",
          12, TINTA, True, align=PP_ALIGN.RIGHT)
    caixa(s, 9.3, 0.52, 3.5, 0.3, "~18 m no perímetro · a confirmar", 8.5,
          SUAVE, align=PP_ALIGN.RIGHT)

    figura(s, f"borda_{b['id']}.jpg", 0.45, 1.2, 7.55)
    caixa(s, 0.45, 5.82, 7.55, 0.25,
          "Piso mantido igual em todas as opções — só a borda muda.",
          9, SUAVE, italico=True)

    x = 8.3
    caixa(s, x, 1.2, 4.55, 0.7, b["resumo"], 12, TINTA, True, entre=1.2)
    for i, (chave, rot) in enumerate(CRIT_BORDA):
        cy = 2.15 + i * 0.28
        caixa(s, x, cy, 1.85, 0.24, rot, 10, TINTA)
        pips(s, x + 1.95, cy + 0.05, b["criterios"][chave],
             largura=0.15, alt=0.1, gap=0.032)

    caixa(s, x, 3.72, 4.55, 0.22, "A FAVOR", 8.5, VERDE, True)
    caixa(s, x, 3.94, 4.55, 0.9, b["a_favor"], 9.5, SUAVE, entre=1.15)
    caixa(s, x, 4.95, 4.55, 0.22, "CONTRA", 8.5, ALERTA, True)
    caixa(s, x, 5.17, 4.55, 0.9, b["contra"], 9.5, SUAVE, entre=1.15)

    if b["alerta"]:
        bloco(s, x, 6.2, 4.55, 0.85, ALERTA)
        caixa(s, x + 0.18, 6.32, 4.2, 0.65, "▲  " + b["alerta"], 9,
              RGBColor(0xFF, 0xFF, 0xFF), True, entre=1.15)


def s_borda_grade(prs, bordas):
    s = novo(prs)
    caixa(s, 0.55, 0.38, 10, 0.45, "As seis bordas lado a lado", 26, TINTA,
          True, TIT)
    caixa(s, 0.55, 0.93, 12.2, 0.3,
          "Mesma vista, mesmo piso, mesma luz. Só a peça de borda muda.",
          11, SUAVE)
    lw, gap = 3.95, 0.19
    for i, b in enumerate(bordas):
        cx = 0.55 + (i % 3) * (lw + gap)
        cy = 1.45 + (i // 3) * 2.82
        figura(s, f"borda_{b['id']}.jpg", cx, cy, lw)
        caixa(s, cx, cy + 2.42, lw, 0.24, f"{b['n']}. {b['nome']}", 10.5,
              TINTA, True)
        caixa(s, cx, cy + 2.66, lw, 0.22,
              f"R$ {b['preco_ml'][0]}–{b['preco_ml'][1]}/m linear", 9, SUAVE)
    rodape(s, "Preços indicativos da peça boleada com pingadeira, instalada, "
              "para a Grande Vitória/ES. A confirmar com marmoraria.")


def s04_como_ler(prs):
    s = novo(prs)
    caixa(s, 0.55, 0.42, 9, 0.45, "Como ler os próximos slides", 26, TINTA,
          True, TIT)
    caixa(s, 0.55, 1.02, 7.4, 0.3,
          "Cada material tem um slide, sempre com o mesmo layout: só o "
          "material muda de um para o outro.", 12, SUAVE)

    caixa(s, 0.55, 1.75, 6, 0.3, "OS CINCO CRITÉRIOS", 10, ACENTO, True)
    explic = [
        ("Antiderrapante", "Aderência com o piso molhado. Critério de segurança."),
        ("Não esquenta", "Conforto ao pé descalço. A calçada não tem sombra."),
        ("Custo baixo", "Material mais mão de obra, por m² instalado."),
        ("Baixa manutenção", "Resistência a mancha, limo, cloro e desbote."),
        ("Resiste à maresia", "Ponta da Fruta é litoral: salinidade cobra caro."),
    ]
    y = 2.2
    for nome, desc in explic:
        caixa(s, 0.55, y, 2.1, 0.3, nome, 11, TINTA, True)
        pips(s, 2.75, y + 0.055, 5)
        caixa(s, 4.05, y, 3.9, 0.3, desc, 9.5, SUAVE)
        y += 0.55

    bloco(s, 0.55, 5.15, 7.4, 0.95, CREME)
    caixa(s, 0.8, 5.35, 6.9, 0.55,
          "Cinco quadrados cheios é sempre o melhor. Em «Custo baixo», "
          "cinco quadrados significa mais barato — não mais caro.",
          11, TINTA, entre=1.25)

    bloco(s, 8.35, 1.75, 4.45, 4.35, CREME)
    caixa(s, 8.65, 2.0, 3.9, 0.3, "DUAS DECISÕES SEPARADAS", 10, ACENTO, True)
    caixa(s, 8.65, 2.42, 3.9, 1.5,
          "Borda e piso não precisam ser do mesmo material. Cada slide marca "
          "se aquele material serve para borda, para piso, ou para os dois.\n\n"
          "A combinação mais comum em obra bem feita é borda de granito com "
          "piso de outro material: a borda pede peça sob medida, com "
          "pingadeira e canto boleado.", 10.5, TINTA, entre=1.25)
    y = 4.35
    for chave in ("sim", "limitado", "nao"):
        rot, cor = SELO[chave]
        bloco(s, 8.65, y, 0.9, 0.26, cor)
        caixa(s, 8.65, y + 0.045, 0.9, 0.26, rot, 8.5,
              RGBColor(0xFF, 0xFF, 0xFF), True, align=PP_ALIGN.CENTER)
        legenda = {"sim": "funciona bem nessa aplicação",
                   "limitado": "funciona, com ressalva",
                   "nao": "não use nessa aplicação"}[chave]
        caixa(s, 9.7, y + 0.03, 2.9, 0.26, legenda, 9.5, SUAVE)
        y += 0.42

    rodape(s, "As montagens são estudo de acabamento, não renderização "
              "fotorrealista. Cor de pedra natural varia por lote: peça "
              "amostra física antes de fechar.")


def s_material(prs, mat, por_id):
    s = novo(prs)
    mid = mat["id"]

    bloco(s, 0, 0, L, 0.92, CREME)
    bloco(s, 0, 0.92, L, 0.02, ACENTO)
    caixa(s, 0.55, 0.16, 0.7, 0.5, f"{mat['n']:02d}", 26, ACENTO, True, TIT)
    caixa(s, 1.32, 0.13, 7.2, 0.4, mat["nome"], 21, TINTA, True, TIT)
    caixa(s, 1.32, 0.575, 7.2, 0.28,
          f"{mat['familia']}  ·  {mat['subtitulo']}", 10.5, SUAVE)

    caixa(s, 9.3, 0.15, 3.5, 0.3,
          f"R$ {mat['preco_m2'][0]}–{mat['preco_m2'][1]} /m² instalado",
          12, TINTA, True, align=PP_ALIGN.RIGHT)
    caixa(s, 9.3, 0.47, 3.5, 0.3, faixa_preco(mat), 14, ACENTO, True,
          align=PP_ALIGN.RIGHT)
    caixa(s, 9.3, 0.72, 3.5, 0.2, "faixa indicativa, confirmar local", 8,
          SUAVE, align=PP_ALIGN.RIGHT)

    figura(s, f"pano_{mid}.jpg", 0.4, 1.18, 6.25)
    figura(s, f"close_{mid}.jpg", 6.82, 1.18, 6.11)
    caixa(s, 0.4, 5.92, 6.25, 0.25, "Vista geral da área", 9, SUAVE, italico=True)
    caixa(s, 6.82, 5.92, 6.11, 0.25, "Detalhe da borda e do encontro com a fibra",
          9, SUAVE, italico=True)

    if mat["borda"] == "nao":
        alt_nome = por_id[BORDA_ALTERNATIVA]["nome"]
        caixa(s, 6.82, 6.15, 6.11, 0.22,
              f"Borda mostrada em {alt_nome} — este material não serve como borda.",
              8, ALERTA, italico=True)

    # Zona inferior: 6,20 a 7,42. Tudo abaixo de 7,45 sai do slide.
    figura(s, f"amostra_{mid}.jpg", 0.4, 6.24, 1.38)
    caixa(s, 0.4, 7.13, 1.38, 0.2, "amostra", 8, SUAVE, align=PP_ALIGN.CENTER)

    x = 2.0
    for i, (chave, rot) in enumerate(CRITERIOS):
        cy = 6.26 + i * 0.205
        caixa(s, x, cy, 1.45, 0.2, rot, 8.5, TINTA)
        pips(s, x + 1.5, cy + 0.035, mat["criterios"][chave],
             largura=0.115, alt=0.085, gap=0.028)

    xs = 4.75
    for rotulo, chave in (("Borda", "borda"), ("Piso", "piso")):
        rot, cor = SELO[mat[chave]]
        caixa(s, xs, 6.28, 0.85, 0.2, rotulo, 8.5, TINTA, True)
        bloco(s, xs, 6.52, 0.88, 0.23, cor)
        caixa(s, xs, 6.558, 0.88, 0.23, rot, 8,
              RGBColor(0xFF, 0xFF, 0xFF), True, align=PP_ALIGN.CENTER)
        xs += 0.98

    yr = 6.42 if mat["borda"] == "nao" else 6.24
    caixa(s, 6.82, yr, 6.11, 0.45, mat["resumo"], 10.5, TINTA, True, entre=1.15)
    caixa(s, 6.82, 6.78, 2.95, 0.2, "A FAVOR", 7.5, VERDE, True)
    caixa(s, 6.82, 6.98, 2.95, 0.5, mat["a_favor"], 7.5, SUAVE, entre=1.08)
    caixa(s, 9.97, 6.78, 2.96, 0.2, "CONTRA", 7.5, ALERTA, True)
    caixa(s, 9.97, 6.98, 2.96, 0.5, mat["contra"], 7.5, SUAVE, entre=1.08)

    # Alerta sobre a propria imagem: nao desloca nada do resto do layout.
    if mat["alerta"]:
        bloco(s, 0.4, 5.35, 6.25, 0.5, ALERTA)
        caixa(s, 0.58, 5.44, 5.9, 0.35, "▲  " + mat["alerta"], 8.5,
              RGBColor(0xFF, 0xFF, 0xFF), True, entre=1.1)


def s15_tabela(prs, mats):
    s = novo(prs)
    caixa(s, 0.55, 0.38, 9, 0.45, "Os dez lado a lado", 26, TINTA, True, TIT)
    caixa(s, 0.55, 0.93, 12.2, 0.3,
          "Ordenado do mais barato para o mais caro. Cinco é sempre o melhor.",
          11, SUAVE)

    cols = ["#", "Material", "Antider-\nrapante", "Não\nesquenta",
            "Custo\nbaixo", "Baixa\nmanut.", "Resiste\nmaresia",
            "R$/m²", "Borda", "Piso"]
    larguras = [0.42, 3.25, 1.02, 0.95, 0.9, 0.95, 1.0, 1.35, 1.05, 1.05]
    ordem = sorted(mats, key=lambda m: sum(m["preco_m2"]) / 2)

    linhas = len(ordem) + 1
    tab = s.shapes.add_table(linhas, len(cols), Inches(0.55), Inches(1.42),
                             Inches(sum(larguras)), Inches(5.2)).table
    for i, w in enumerate(larguras):
        tab.columns[i].width = Inches(w)
    tab.rows[0].height = Inches(0.52)

    for j, c in enumerate(cols):
        cel = tab.cell(0, j)
        cel.text = c
        cel.fill.solid()
        cel.fill.fore_color.rgb = RGBColor(0x12, 0x2A, 0x33)
        cel.vertical_anchor = MSO_ANCHOR.MIDDLE
        cel.margin_left = cel.margin_right = Inches(0.05)
        for p in cel.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(9)
                r.font.bold = True
                r.font.name = TXT
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i, m in enumerate(ordem, start=1):
        tab.rows[i].height = Inches(0.42)
        vals = [str(m["n"]), m["nome"]] + \
               [str(m["criterios"][k]) for k, _ in CRITERIOS] + \
               [f"{m['preco_m2'][0]}–{m['preco_m2'][1]}",
                SELO[m["borda"]][0], SELO[m["piso"]][0]]
        for j, v in enumerate(vals):
            cel = tab.cell(i, j)
            cel.text = v
            cel.fill.solid()
            cel.fill.fore_color.rgb = (RGBColor(0xFF, 0xFF, 0xFF)
                                       if i % 2 else CREME)
            cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            cel.margin_left = cel.margin_right = Inches(0.05)
            for p in cel.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(9.5)
                    r.font.name = TXT
                    r.font.bold = j in (0, 1)
                    if j in (8, 9):
                        r.font.color.rgb = SELO[
                            m["borda" if j == 8 else "piso"]][1]
                    elif j >= 2 and j <= 6:
                        nota = m["criterios"][CRITERIOS[j - 2][0]]
                        r.font.color.rgb = (ALERTA if nota <= 2
                                            else TINTA if nota <= 3 else VERDE)
                    else:
                        r.font.color.rgb = TINTA

    rodape(s, "Preços indicativos para a Grande Vitória/ES, a confirmar com "
              "fornecedor local. Área estimada por foto: 60 a 75 m² de piso "
              "e ~18 m de borda.")


def s16_comparador(prs, mats):
    s = novo(prs)
    caixa(s, 0.55, 0.38, 9, 0.45, "Os dez na mesma foto", 26, TINTA, True, TIT)
    caixa(s, 0.55, 0.93, 12.2, 0.3,
          "Mesma vista, mesma luz, mesmo enquadramento: só o material muda.",
          11, SUAVE)

    lw, gap = 2.34, 0.13         # 5 colunas dentro das margens de 0,55
    for i, m in enumerate(mats):
        cx = 0.55 + (i % 5) * (lw + gap)
        cy = 1.5 + (i // 5) * 2.72
        figura(s, f"mini_{m['id']}.jpg", cx, cy, lw)
        caixa(s, cx, cy + 1.86, lw, 0.22, f"{m['n']:02d}  {m['nome']}", 8.5,
              TINTA, True, entre=1.05)
        caixa(s, cx, cy + 2.12, lw, 0.2,
              f"R$ {m['preco_m2'][0]}–{m['preco_m2'][1]}/m²", 8, SUAVE)

    rodape(s, "Estudo de acabamento sobre foto real. Cor de pedra natural "
              "varia por lote — peça amostra física antes de fechar.")


def s17_recomendacao(prs, por_id):
    s = novo(prs)
    caixa(s, 0.55, 0.38, 9, 0.45, "Três combinações que eu levaria adiante",
          26, TINTA, True, TIT)
    caixa(s, 0.55, 0.93, 12.2, 0.3,
          "Borda e piso escolhidos juntos. Todas com borda em meia-cana.",
          11, SUAVE)

    porques = {
        "equilibrio":
            "O meio-termo defensável. Claro o bastante para não ferver ao "
            "meio-dia, resistente o bastante para não pedir "
            "impermeabilização, e barato pela proximidade de Cachoeiro. "
            "Peça única para borda e piso simplifica a compra e a obra.",
        "conforto":
            "Se o pé descalço no sol das duas da tarde é a sua prioridade, "
            "é esta. O quartzito é o mais frio da lista; a borda em granito "
            "assume o lugar onde mais se pisa molhado e mais se arrasta o "
            "corpo, que é onde o quartzito sofreria.",
        "manutencao":
            "Para quem não quer pensar na área de novo. Porcelanato não "
            "mancha, não cria limo e ignora a maresia; a borda em granito "
            "resolve o recorte da peça sob medida, que porcelanato de "
            "90×90 não entrega bem no perímetro.",
    }
    ressalvas = {
        "equilibrio": "Granito claro flameado mostra pó: pede lavagem mais frequente.",
        "conforto": "Quartzito é poroso — impermeabilização é obrigatória, não opcional.",
        "manutencao": "Exige contrapiso impecável; peça grande trinca se a base ceder.",
    }

    cw = 4.06
    for i, (slug, id_piso, id_borda, rotulo) in enumerate(COMBOS):
        cx = 0.55 + i * (cw + 0.09)
        bloco(s, cx, 1.42, cw, 5.55, CREME)
        bloco(s, cx, 1.42, cw, 0.045, ACENTO)
        caixa(s, cx + 0.22, 1.62, cw - 0.44, 0.32, rotulo, 15, TINTA, True, TIT)
        figura(s, f"combo_{slug}.jpg", cx + 0.22, 2.05, cw - 0.44)
        caixa(s, cx + 0.22, 4.98, cw - 0.44, 0.5,
              f"Piso: {por_id[id_piso]['nome']}\nBorda: {por_id[id_borda]['nome']}",
              9.5, ACENTO, True, entre=1.2)
        caixa(s, cx + 0.22, 5.55, cw - 0.44, 1.0, porques[slug], 9, SUAVE,
              entre=1.15)
        caixa(s, cx + 0.22, 6.62, cw - 0.44, 0.3, "Ressalva: " + ressalvas[slug],
              8.5, ALERTA, entre=1.1)

    rodape(s, "Recomendação de partida, não veredito. O critério que mais "
              "pesa é seu: quem usa a área descalço no meio do dia decide "
              "diferente de quem usa no fim da tarde.")


def s18_proximos(prs):
    s = novo(prs)
    caixa(s, 0.55, 0.38, 9, 0.45, "Próximos passos", 26, TINTA, True, TIT)
    caixa(s, 0.55, 0.93, 12.2, 0.3,
          "O que confirmar antes de pedir orçamento, e o que exigir de quem "
          "for executar.", 11, SUAVE)

    colunas = [
        ("CONFIRMAR NO LOCAL", ACENTO, [
            "Medir a piscina com trena: comprimento, largura e perímetro reais.",
            "Medir a calçada nos quatro lados — a largura varia.",
            "Conferir se há caimento hoje e para onde a água corre.",
            "Verificar o estado do contrapiso sob a São Tomé ao remover "
            "a primeira placa.",
            "Fotografar o encontro da fibra com o concreto nos quatro cantos.",
        ]),
        ("PERGUNTAR AO MARMORISTA", VERDE, [
            "O material tem placa em estoque ou é sob encomenda? Qual o prazo?",
            "Consegue peça de borda sob medida, boleada e com pingadeira?",
            "Qual a espessura da peça de borda e da placa de piso?",
            "Pode fornecer amostra física de 20×20 do lote que será usado?",
            "Qual a variação de tom esperada entre chapas do mesmo lote?",
        ]),
        ("EXIGIR DE QUEM EXECUTAR", ALERTA, [
            "Junta flexível em mastique PU entre a borda e a fibra — "
            "nunca rejunte rígido.",
            "Caimento de 1 a 1,5% para fora da piscina, conferido com nível.",
            "Rejunte epóxi em toda a área, não cimentício.",
            "Juntas de dilatação a cada ~3 m.",
            "Toda ferragem e parafuso em inox 316 (A4).",
        ]),
    ]

    cw = 4.06
    for i, (titulo, cor, itens) in enumerate(colunas):
        cx = 0.55 + i * (cw + 0.09)
        bloco(s, cx, 1.45, cw, 5.1, CREME)
        bloco(s, cx, 1.45, cw, 0.045, cor)
        caixa(s, cx + 0.25, 1.68, cw - 0.5, 0.3, titulo, 10.5, cor, True)
        y = 2.12
        for item in itens:
            bloco(s, cx + 0.25, y + 0.055, 0.115, 0.115, cor)
            caixa(s, cx + 0.5, y, cw - 0.75, 0.8, item, 10, TINTA, entre=1.18)
            y += 0.88

    bloco(s, 0.55, 6.72, 12.23, 0.5, RGBColor(0x12, 0x2A, 0x33))
    caixa(s, 0.85, 6.86, 11.6, 0.3,
          "Antes de fechar qualquer material: peça amostra física e veja a "
          "placa sob o sol do local, molhada e seca.",
          11, RGBColor(0xFF, 0xFF, 0xFF), True)


# ---------------------------------------------------------------- principal


def main():
    mats = M.carrega_materiais()
    por_id = {m["id"]: m for m in mats}

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(L), Inches(A)

    bordas = M.carrega_bordas()

    s01_capa(prs)
    s02_diagnostico(prs)
    s03_o_que_muda(prs)
    # A borda vem antes do piso: e o defeito principal do diagnostico.
    s_borda_intro(prs)
    for b in bordas:
        s_borda(prs, b)
    s_borda_grade(prs, bordas)
    s04_como_ler(prs)
    for m in mats:
        s_material(prs, m, por_id)
    s15_tabela(prs, mats)
    s16_comparador(prs, mats)
    s17_recomendacao(prs, por_id)
    s18_proximos(prs)

    destino = salva(prs, DESTINO)
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides")
    print(destino)


if __name__ == "__main__":
    main()
